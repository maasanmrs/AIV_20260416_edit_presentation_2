#!/usr/bin/env python3
"""
Corporate pitch PPTX generator — AIVALIX branded.

Features:
- AIVALIX black/white base + client accent color
- Dual logo placement (AIVALIX + client) on every slide
- One-slide-at-a-time pipeline with checkpoint saves
- Slide types: cover, agenda, section, content_bullets, content_table,
               two_column, team, back_cover
- Auto-shrink text, crop-to-fill images, dark overlay on background images
- Two-pass image workflow: emit image slots → generate images → insert
"""
import argparse, json, os, io, sys, math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Dimensions ──────────────────────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

FONT_JP = "Noto Sans JP"
FONT_EN = "Noto Sans"

# ── AIVALIX Color Scheme ────────────────────────────────────────────────
C = {
    "black":         RGBColor(0x00, 0x00, 0x00),
    "white":         RGBColor(0xFF, 0xFF, 0xFF),
    "dark_bg":       RGBColor(0x0A, 0x0A, 0x0A),
    "light_bg":      RGBColor(0xF5, 0xF5, 0xF5),
    "text_dark":     RGBColor(0x1A, 0x1A, 0x1A),
    "text_light":    RGBColor(0x99, 0x99, 0x99),
    "sep":           RGBColor(0xE0, 0xE0, 0xE0),
    "accent":        RGBColor(0x1A, 0x36, 0x5D),  # Default; overridden by client color
    "accent_light":  RGBColor(0x7A, 0xB8, 0xE0),  # Light version for use on dark backgrounds
}


def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ── Shape Helpers ───────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_textbox(slide, text, x, y, w, h, font_size, bold=False,
                color=None, h_align=PP_ALIGN.LEFT, font=None,
                word_wrap=True, auto_shrink=False, min_size=10,
                v_align=MSO_ANCHOR.MIDDLE):
    color = color or C["white"]
    font = font or FONT_JP

    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = v_align

    p = tf.paragraphs[0]
    p.alignment = h_align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font

    if auto_shrink:
        for fs in [font_size, font_size - 1, font_size - 2, font_size - 3, min_size]:
            run.font.size = Pt(fs)
            if fs <= min_size:
                break
            approx_cpl = w / Pt(fs) * 1.5
            if len(text) <= approx_cpl:
                break
    else:
        run.font.size = Pt(font_size)
    return txBox


# ── Image Helpers ───────────────────────────────────────────────────────

def crop_image_to_box(img_path, target_w_emu, target_h_emu):
    try:
        from PIL import Image as PI
    except ImportError:
        return None
    if not img_path or not os.path.exists(img_path):
        return None
    img = PI.open(img_path).convert("RGB")
    sw, sh = img.size
    tr = target_w_emu / target_h_emu
    sr = sw / sh
    if sr > tr:
        nw = int(sh * tr)
        img = img.crop(((sw - nw) // 2, 0, (sw - nw) // 2 + nw, sh))
    else:
        nh = int(sw / tr)
        img = img.crop((0, (sh - nh) // 2, sw, (sh - nh) // 2 + nh))
    buf = io.BytesIO()
    img.save(buf, "JPEG", quality=92)
    buf.seek(0)
    return buf


def place_image(slide, img_path, x, y, w, h):
    if not img_path or not os.path.exists(img_path):
        return False
    buf = crop_image_to_box(img_path, w, h)
    if buf:
        slide.shapes.add_picture(buf, x, y, w, h)
        return True
    try:
        slide.shapes.add_picture(img_path, x, y, w, h)
        return True
    except Exception:
        return False


def place_logo(slide, logo_path, x, y, max_w, max_h):
    """Place logo preserving aspect ratio within max bounds."""
    if not logo_path or not os.path.exists(logo_path):
        return False
    try:
        from PIL import Image as PI
        img = PI.open(logo_path)
        iw, ih = img.size
        ratio = min(max_w / Inches(iw / 96), max_h / Inches(ih / 96))
        w = Inches(iw / 96) * ratio if ratio < 1 else Inches(iw / 96)
        h = Inches(ih / 96) * ratio if ratio < 1 else Inches(ih / 96)
        # Clamp to max
        if w > max_w:
            h = h * max_w / w
            w = max_w
        if h > max_h:
            w = w * max_h / h
            h = max_h
        slide.shapes.add_picture(logo_path, int(x), int(y), int(w), int(h))
        return True
    except Exception:
        try:
            slide.shapes.add_picture(logo_path, int(x), int(y), int(max_w), int(max_h))
            return True
        except Exception:
            return False


def add_dark_overlay(slide, x, y, w, h, alpha=150):
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    ov = slide.shapes.add_shape(1, x, y, w, h)
    ov.fill.solid()
    ov.fill.fore_color.rgb = C["black"]
    ov.line.fill.background()
    solidFill = ov._element.find(".//" + qn("a:solidFill"))
    if solidFill is not None:
        srgb = solidFill.find(qn("a:srgbClr"))
        if srgb is not None:
            ae = etree.SubElement(srgb, qn("a:alpha"))
            ae.set("val", str(int((1 - alpha / 255) * 100000)))


def grad_placeholder(slide, x, y, w, h, colors):
    add_rect(slide, x, y, w, h, colors["dark_bg"])
    band = h // 6
    for i, f in enumerate([0.08, 0.12, 0.18, 0.22, 0.28, 0.35]):
        shade = RGBColor(
            min(255, int(colors["accent"][0] * f)),
            min(255, int(colors["accent"][1] * f)),
            min(255, int(colors["accent"][2] * f + 0x20 * (1 - f))),
        )
        add_rect(slide, x, y + band * i, w, band, shade)


# ── Image Slot Tracking ────────────────────────────────────────────────
# Collects exact pixel dimensions for each image area during PPTX build.
# Used in two-pass workflow: 1) build PPTX skeleton → 2) generate images → 3) insert.

_IMAGE_SLOTS = []   # list of dicts: {slide_idx, x, y, w, h, w_px, h_px, style_hint, prompt_hint}

DPI = 96  # PowerPoint standard DPI

def _emu_to_px(emu):
    """Convert EMU to pixels at 96 DPI."""
    return round(emu * DPI / 914400)

def _record_image_slot(slide_idx, x, y, w, h, style_hint="corporate", prompt_hint=""):
    """Record an image slot with exact pixel dimensions for Nano Banana Pro."""
    _IMAGE_SLOTS.append({
        "slide_idx": slide_idx,
        "x_emu": int(x), "y_emu": int(y),
        "w_emu": int(w), "h_emu": int(h),
        "w_px": _emu_to_px(int(w)),
        "h_px": _emu_to_px(int(h)),
        "style_hint": style_hint,
        "prompt_hint": prompt_hint,
    })


# ── Logo Placement ──────────────────────────────────────────────────────

def place_logos(slide, aivalix_logo, client_logo, slide_type="content"):
    """Place both logos on the slide."""
    margin = Inches(0.3)
    if slide_type == "cover":
        # Cover: logos at top
        logo_h = Inches(0.45)
        logo_w = Inches(1.8)
        place_logo(slide, aivalix_logo, margin, margin, logo_w, logo_h)
        if client_logo:
            place_logo(slide, client_logo,
                        SLIDE_W - logo_w - margin, margin, logo_w, logo_h)
    elif slide_type == "back_cover":
        # Back cover: logos centered at bottom
        logo_h = Inches(0.5)
        logo_w = Inches(2.0)
        center_x = (SLIDE_W - logo_w * 2 - Inches(1.0)) // 2
        place_logo(slide, aivalix_logo, center_x, SLIDE_H - logo_h - Inches(0.8),
                    logo_w, logo_h)
        if client_logo:
            place_logo(slide, client_logo,
                        center_x + logo_w + Inches(1.0),
                        SLIDE_H - logo_h - Inches(0.8), logo_w, logo_h)
    else:
        # Content slides: AIVALIX logo is placed in header by _draw_content_header.
        # Only place client (TGES) logo at bottom-left.
        if client_logo:
            logo_h = Inches(0.28)
            logo_w = Inches(1.2)
            y = SLIDE_H - logo_h - margin
            place_logo(slide, client_logo, margin, y, logo_w, logo_h)


# ── Slide Makers ────────────────────────────────────────────────────────

def make_cover(prs, data, img_path, colors, logos, *, slide_idx=0):
    """Cover slide: full-bleed image, dark overlay, title, logos."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, colors["black"])

    # Record image slot (full-bleed background)
    _record_image_slot(slide_idx, 0, 0, SLIDE_W, SLIDE_H,
                       style_hint="corporate",
                       prompt_hint=data.get("image_prompt", ""))

    placed = place_image(slide, img_path, 0, 0, SLIDE_W, SLIDE_H)
    if not placed:
        grad_placeholder(slide, 0, 0, SLIDE_W, SLIDE_H, colors)
    add_dark_overlay(slide, 0, 0, SLIDE_W, SLIDE_H, alpha=145)

    cx = Inches(0.55)
    cw = Inches(8.5)

    # Title
    add_textbox(slide, data.get("title", ""),
                cx, Inches(2.0), cw, Inches(2.5),
                48, bold=True, color=colors["white"], font=FONT_JP,
                word_wrap=True, v_align=MSO_ANCHOR.MIDDLE)

    # Subtitle
    add_textbox(slide, data.get("subtitle", ""),
                cx, Inches(4.6), cw, Inches(1.0),
                18, color=colors["text_light"], font=FONT_JP,
                word_wrap=True, v_align=MSO_ANCHOR.MIDDLE)

    # Date
    date_str = data.get("date", "")
    add_textbox(slide, date_str,
                cx, SLIDE_H - Inches(1.0), Inches(3), Inches(0.5),
                13, color=colors["text_light"], font=FONT_EN,
                v_align=MSO_ANCHOR.MIDDLE)

    # Accent line (use light version for contrast on dark bg)
    add_rect(slide, cx, Inches(1.85), Inches(4.5), Inches(0.04),
             colors.get("accent_light", colors["accent"]))

    place_logos(slide, logos["aivalix"], logos["client"], "cover")


def make_agenda(prs, data, colors, logos):
    """Agenda/TOC slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, colors["white"])

    # Title bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.4), colors["black"])
    add_textbox(slide, data.get("title", "Agenda"),
                Inches(0.55), Inches(0.3), Inches(8), Inches(0.8),
                32, bold=True, color=colors["white"], font=FONT_JP,
                v_align=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 0, Inches(1.4), SLIDE_W, Inches(0.05), colors["accent"])

    items = data.get("items", [])
    n = len(items)
    start_y = Inches(1.8)
    avail_h = SLIDE_H - start_y - Inches(0.8)
    row_h = min(Inches(0.75), avail_h / max(n, 1))

    for i, item in enumerate(items):
        y = start_y + i * row_h
        # Number circle
        num_size = Inches(0.4)
        num_x = Inches(0.8)
        num_y = y + (row_h - num_size) / 2
        add_rect(slide, num_x, num_y, num_size, num_size, colors["accent"])
        add_textbox(slide, str(i + 1),
                    num_x, num_y, num_size, num_size,
                    16, bold=True, color=colors["white"], font=FONT_EN,
                    h_align=PP_ALIGN.CENTER, v_align=MSO_ANCHOR.MIDDLE)
        # Item text
        add_textbox(slide, item,
                    num_x + num_size + Inches(0.3), y, Inches(9), row_h,
                    18, color=colors["text_dark"], font=FONT_JP,
                    v_align=MSO_ANCHOR.MIDDLE)
        # Separator
        if i < n - 1:
            add_rect(slide, Inches(0.8), y + row_h - Inches(0.01),
                     Inches(10), Inches(0.01), colors["sep"])

    place_logos(slide, logos["aivalix"], logos["client"], "content")


def make_section(prs, data, img_path, colors, logos, *, slide_idx=0):
    """Section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, colors["black"])

    # Record image slot (full-bleed background)
    _record_image_slot(slide_idx, 0, 0, SLIDE_W, SLIDE_H,
                       style_hint="abstract",
                       prompt_hint=data.get("image_prompt", ""))

    placed = place_image(slide, img_path, 0, 0, SLIDE_W, SLIDE_H)
    if not placed:
        grad_placeholder(slide, 0, 0, SLIDE_W, SLIDE_H, colors)
    add_dark_overlay(slide, 0, 0, SLIDE_W, SLIDE_H, alpha=165)

    cx = Inches(0.6)
    cw = SLIDE_W - Inches(1.2)

    # Large section number (decorative, subtle)
    num = data.get("number", 1)
    add_textbox(slide, f"0{num}" if num < 10 else str(num),
                Inches(0.2), Inches(0.4), Inches(4.5), Inches(4.2),
                140, bold=True,
                color=RGBColor(0x15, 0x20, 0x30), font=FONT_EN,
                v_align=MSO_ANCHOR.MIDDLE)

    add_rect(slide, cx, Inches(2.9), Inches(1.5), Inches(0.06), colors["accent"])

    add_textbox(slide, data.get("title", ""),
                cx, Inches(3.05), cw, Inches(1.5),
                50, bold=True, color=colors["white"], font=FONT_JP,
                word_wrap=True, auto_shrink=True, min_size=32,
                v_align=MSO_ANCHOR.MIDDLE)

    add_textbox(slide, data.get("subtitle", ""),
                cx, Inches(4.7), cw, Inches(0.95),
                18, color=colors["text_light"], font=FONT_JP,
                word_wrap=True, auto_shrink=True, min_size=13,
                v_align=MSO_ANCHOR.MIDDLE)

    place_logos(slide, logos["aivalix"], logos["client"], "content")


def _draw_content_header(slide, data, colors, logos=None):
    """Draw compact header bar (single-line title) + prominent key message.

    Layout (top to bottom):
    1. Black header bar (~1.0"): section label + title (MUST be 1 line) + AIVALIX logo
    2. Accent line (0.05")
    3. Key message (full width, 20pt bold black, 1-2 lines) + separator
    Returns content_start_y — images and body content go BELOW this line.
    """
    header_h = Inches(1.0)  # Compact: single-line title only

    add_rect(slide, 0, 0, SLIDE_W, header_h, colors["black"])
    sec_num = data.get("section_number", 1)
    sec_title = data.get("section_title", "").upper()
    label_color = colors.get("accent_light", colors["white"])
    logo_reserve = Inches(3.0)
    text_x = Inches(0.5)
    full_w = SLIDE_W - Inches(1.0)

    add_textbox(slide, f"0{sec_num}  {sec_title}" if sec_num < 10 else f"{sec_num}  {sec_title}",
                text_x, Inches(0.08), full_w - logo_reserve, Inches(0.30),
                12, color=label_color, font=FONT_EN,
                v_align=MSO_ANCHOR.MIDDLE)
    # Title — SINGLE LINE, aggressive auto_shrink, word_wrap=False
    add_textbox(slide, data.get("title", ""),
                text_x, Inches(0.36), full_w - logo_reserve, Inches(0.52),
                26, bold=True, color=colors["white"], font=FONT_JP,
                word_wrap=False, auto_shrink=True, min_size=16,
                v_align=MSO_ANCHOR.MIDDLE)

    add_rect(slide, 0, header_h, SLIDE_W, Inches(0.05), colors["accent"])

    # AIVALIX logo — 2.5x title text height
    if logos and logos.get("aivalix"):
        logo_h = Inches(0.90)
        logo_w = Inches(2.62)
        logo_x = SLIDE_W - logo_w - Inches(0.15)
        logo_y = Inches(0.05)
        place_logo(slide, logos["aivalix"], logo_x, logo_y, logo_w, logo_h)

    content_y = header_h + Inches(0.05)

    # Key message — FULL WIDTH, 20pt bold black, 1-2 lines
    # Size is between body text (13pt) and title (26pt)
    key_msg = data.get("key_message", "")
    if key_msg:
        km_y = content_y + Inches(0.10)
        km_h = Inches(0.72)  # Room for 1-2 lines at 20pt
        add_textbox(slide, key_msg,
                    text_x, km_y, full_w, km_h,
                    20, bold=True, color=C["text_dark"],
                    font=FONT_JP, word_wrap=True, auto_shrink=True, min_size=16,
                    v_align=MSO_ANCHOR.MIDDLE)
        sep_y = km_y + km_h + Inches(0.04)
        add_rect(slide, text_x, sep_y, full_w, Inches(0.015), C["sep"])
        content_y = sep_y + Inches(0.08)

    return content_y


def _content_layout_below_header(data, img_path, content_start_y):
    """Calculate image + text layout for the content area BELOW header.

    Images are placed only in the content area (below header + key message),
    never in the black header bar. Image dimensions adapt to remaining space.
    Returns: placement, img_x, img_y, img_w, img_h, text_x, text_w
    """
    placement = data.get("image_placement", "auto")
    has_img = bool(img_path and os.path.exists(img_path))
    if placement == "auto":
        placement = "right" if has_img else "none"
    if not has_img and placement != "none":
        placement = "none"

    footer_margin = Inches(0.55)
    content_h = SLIDE_H - content_start_y - footer_margin
    text_x = Inches(0.5)

    if placement == "right":
        # Adaptive image width based on content density
        n_items = len(data.get("content", data.get("table_rows", [])))
        if n_items <= 3:
            img_w = Inches(5.0)   # Few items → wider image
        elif n_items <= 5:
            img_w = Inches(4.5)   # Medium → balanced
        else:
            img_w = Inches(4.0)   # Many items → more text space
        img_h = content_h
        img_x = SLIDE_W - img_w
        img_y = content_start_y
        text_w = img_x - text_x - Inches(0.2)
    elif placement == "background":
        img_w = SLIDE_W
        img_h = SLIDE_H
        img_x = 0
        img_y = 0
        text_w = SLIDE_W - Inches(1.0)
    else:
        img_w = 0
        img_h = 0
        img_x = 0
        img_y = 0
        text_w = SLIDE_W - Inches(1.0)

    return placement, img_x, img_y, img_w, img_h, text_x, text_w


def make_content_bullets(prs, data, img_path, colors, logos, *, slide_idx=0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, colors["light_bg"])

    # 1. Draw header + key message FIRST (full width, no image)
    content_start_y = _draw_content_header(slide, data, colors, logos)

    # 2. Calculate image layout in content area BELOW header
    placement, img_x, img_y, img_w, img_h, text_x, text_w = \
        _content_layout_below_header(data, img_path, content_start_y)

    # Record image slot (below header only)
    img_placement = data.get("image_placement", "auto")
    if img_placement in ("right", "auto"):
        n_items = len(data.get("content", []))
        slot_w = Inches(5.0) if n_items <= 3 else (Inches(4.5) if n_items <= 5 else Inches(4.0))
        slot_x = SLIDE_W - slot_w
        slot_y = content_start_y
        slot_h = SLIDE_H - content_start_y - Inches(0.55)
        _record_image_slot(slide_idx, slot_x, slot_y, slot_w, slot_h,
                           style_hint=data.get("image_style", "diagram"),
                           prompt_hint=data.get("image_prompt", ""))
    elif img_placement == "background":
        _record_image_slot(slide_idx, 0, 0, SLIDE_W, SLIDE_H,
                           style_hint=data.get("image_style", "corporate"),
                           prompt_hint=data.get("image_prompt", ""))

    # 3. Place image in content area only (below header)
    if placement == "right" and img_w > 0:
        placed = place_image(slide, img_path, img_x, img_y, img_w, img_h)
        if not placed:
            grad_placeholder(slide, img_x, img_y, img_w, img_h, colors)
        add_rect(slide, img_x, img_y, Inches(0.05), img_h, colors["accent"])
    elif placement == "background" and img_w > 0:
        placed = place_image(slide, img_path, 0, 0, SLIDE_W, SLIDE_H)
        if placed:
            add_dark_overlay(slide, 0, content_start_y, SLIDE_W,
                             SLIDE_H - content_start_y, alpha=60)

    # 4. Place bullets in remaining text area
    bullets = data.get("content", [])
    n = len(bullets)
    avail_h = SLIDE_H - content_start_y - Inches(0.7)
    row_h = min(Inches(0.98), avail_h / max(n, 1))
    base_fs = 15.5 if n <= 5 else (13.5 if n <= 7 else 12.0)

    for i, bullet in enumerate(bullets):
        y = content_start_y + Inches(0.1) + i * row_h
        dot_size = Inches(0.11)
        dot_x = text_x + Inches(0.07)
        dot_y = y + (row_h / 2) - (dot_size / 2)
        add_rect(slide, dot_x, dot_y, dot_size, dot_size, colors["accent"])

        tx = dot_x + dot_size + Inches(0.14)
        tw = text_w - dot_size - Inches(0.21)
        th = row_h

        tbx = slide.shapes.add_textbox(tx, y, tw, th)
        tf = tbx.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = bullet
        run.font.name = FONT_JP
        run.font.color.rgb = C["text_dark"]
        chosen = base_fs
        for fs in [base_fs, base_fs - 1, base_fs - 2, 12.5, 12.0, 11.5, 11.0]:
            est_cpl = tw / Pt(fs) * 1.5
            chosen = fs
            if len(bullet) / max(est_cpl, 1) < 2.1 or fs <= 11.0:
                break
        run.font.size = Pt(chosen)

        if i < n - 1:
            add_rect(slide, text_x, y + row_h - Inches(0.015),
                     text_w, Inches(0.015), C["sep"])

    place_logos(slide, logos["aivalix"], logos["client"], "content")


def make_content_table(prs, data, img_path, colors, logos, *, slide_idx=0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, colors["light_bg"])

    # 1. Draw header + key message FIRST (full width)
    content_start_y = _draw_content_header(slide, data, colors, logos)

    # 2. Calculate image layout in content area BELOW header
    placement, img_x, img_y, img_w, img_h, text_x, text_w = \
        _content_layout_below_header(data, img_path, content_start_y)

    # Record image slot (below header only)
    img_placement = data.get("image_placement", "auto")
    if img_placement in ("right", "auto") and img_placement != "none":
        n_items = len(data.get("table_rows", []))
        slot_w = Inches(5.0) if n_items <= 3 else (Inches(4.5) if n_items <= 5 else Inches(4.0))
        slot_x = SLIDE_W - slot_w
        slot_y = content_start_y
        slot_h = SLIDE_H - content_start_y - Inches(0.55)
        _record_image_slot(slide_idx, slot_x, slot_y, slot_w, slot_h,
                           style_hint=data.get("image_style", "diagram"),
                           prompt_hint=data.get("image_prompt", ""))
    elif img_placement == "background":
        _record_image_slot(slide_idx, 0, 0, SLIDE_W, SLIDE_H,
                           style_hint=data.get("image_style", "corporate"),
                           prompt_hint=data.get("image_prompt", ""))

    # 3. Place image in content area only
    if placement == "right" and img_w > 0:
        placed = place_image(slide, img_path, img_x, img_y, img_w, img_h)
        if not placed:
            grad_placeholder(slide, img_x, img_y, img_w, img_h, colors)
        add_rect(slide, img_x, img_y, Inches(0.05), img_h, colors["accent"])
    elif placement == "background" and img_w > 0:
        placed = place_image(slide, img_path, 0, 0, SLIDE_W, SLIDE_H)
        if placed:
            add_dark_overlay(slide, 0, content_start_y, SLIDE_W,
                             SLIDE_H - content_start_y, alpha=60)

    headers = data.get("table_headers", [])
    rows = data.get("table_rows", [])
    n_cols = len(headers)
    n_rows = len(rows)
    if n_cols == 0:
        place_logos(slide, logos["aivalix"], logos["client"], "content")
        return

    raw = data.get("table_col_widths")
    if raw and len(raw) == n_cols:
        total = sum(raw)
        col_widths = [text_w * (r / total) for r in raw]
    elif n_cols == 2:
        col_widths = [text_w * 0.28, text_w * 0.72]
    elif n_cols == 3:
        col_widths = [text_w * 0.20, text_w * 0.35, text_w * 0.45]
    elif n_cols == 4:
        col_widths = [text_w * 0.15, text_w * 0.25, text_w * 0.35, text_w * 0.25]
    else:
        col_widths = [text_w / n_cols] * n_cols

    hdr_h = Inches(0.50)
    avail = SLIDE_H - content_start_y - hdr_h - Inches(0.7)  # Leave room for footer logo

    # --- Adaptive row height & font size based on content density ---
    # Measure max cell text length per row to decide how tall each row needs to be
    max_cell_len = max((max((len(str(c)) for c in row[:n_cols]), default=0) for row in rows), default=20)
    # Estimate column width in characters (for the widest column)
    widest_col_emu = max(col_widths) if col_widths else text_w
    # Font size: larger when fewer rows & short text, smaller when dense
    if n_rows <= 4 and max_cell_len <= 30:
        cell_fs = 14.0
    elif n_rows <= 5:
        cell_fs = 13.0
    elif n_rows <= 7:
        cell_fs = 12.0
    else:
        cell_fs = 11.0
    # Row height: use available space but shrink rows when text is short
    chars_per_line = widest_col_emu / Pt(cell_fs) * 1.5
    est_lines = max(1, max_cell_len / max(chars_per_line, 1))
    min_row_h = Inches(0.42) if est_lines <= 1 else Inches(0.55) if est_lines <= 2 else Inches(0.72)
    row_h = max(min_row_h, min(Inches(1.05), avail / max(n_rows, 1)))

    x0, y0 = text_x, content_start_y + Inches(0.05)

    x = x0
    for j, hdr in enumerate(headers):
        add_rect(slide, x, y0, col_widths[j], hdr_h, colors["accent"])
        add_textbox(slide, hdr,
                    x + Inches(0.07), y0, col_widths[j] - Inches(0.1), hdr_h,
                    12, bold=True, color=colors["white"], font=FONT_JP,
                    word_wrap=True, auto_shrink=True, min_size=9,
                    v_align=MSO_ANCHOR.MIDDLE)
        x += col_widths[j]

    for i, row in enumerate(rows):
        y = y0 + hdr_h + row_h * i
        if y + row_h > SLIDE_H - Inches(0.05):
            break
        bg = C["white"] if i % 2 == 0 else colors["light_bg"]
        x = x0
        for j, cell in enumerate(row[:n_cols]):
            add_rect(slide, x, y, col_widths[j], row_h, bg)
            txt_c = colors["accent"] if j == 0 else C["text_dark"]
            # Enable word wrap for all cells; top-align when wrapping
            add_textbox(slide, str(cell).replace("\\n", "\n"),
                        x + Inches(0.07), y + Inches(0.08),
                        col_widths[j] - Inches(0.14), row_h - Inches(0.16),
                        cell_fs, bold=(j == 0), color=txt_c, font=FONT_JP,
                        word_wrap=True, auto_shrink=True, min_size=9,
                        v_align=MSO_ANCHOR.TOP)
            x += col_widths[j]
        add_rect(slide, x0, y + row_h - Inches(0.015),
                 sum(col_widths), Inches(0.015), C["sep"])

    place_logos(slide, logos["aivalix"], logos["client"], "content")


def make_two_column(prs, data, img_path, colors, logos):
    """Two-column comparison slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, colors["light_bg"])

    content_start_y = _draw_content_header(slide, data, colors, logos)

    col_w = (SLIDE_W - Inches(1.5)) / 2
    left_data = data.get("left", {})
    right_data = data.get("right", {})

    col_hdr_y = content_start_y + Inches(0.1)
    col_body_y = col_hdr_y + Inches(0.65)

    for col_idx, (col_data, col_x) in enumerate([
        (left_data, Inches(0.5)),
        (right_data, Inches(0.5) + col_w + Inches(0.5)),
    ]):
        # Column header
        add_rect(slide, col_x, col_hdr_y, col_w, Inches(0.55), colors["accent"])
        add_textbox(slide, col_data.get("title", ""),
                    col_x + Inches(0.1), col_hdr_y, col_w - Inches(0.2), Inches(0.55),
                    16, bold=True, color=colors["white"], font=FONT_JP,
                    h_align=PP_ALIGN.CENTER, v_align=MSO_ANCHOR.MIDDLE)

        items = col_data.get("items", [])
        item_avail = SLIDE_H - col_body_y - Inches(0.7)
        item_h = min(Inches(0.65), item_avail / max(len(items), 1))
        for i, item in enumerate(items):
            y = col_body_y + i * item_h
            dot_size = Inches(0.09)
            add_rect(slide, col_x + Inches(0.1), y + Inches(0.22),
                     dot_size, dot_size, colors["accent"])
            add_textbox(slide, item,
                        col_x + Inches(0.3), y, col_w - Inches(0.4), Inches(0.6),
                        13, color=colors["text_dark"], font=FONT_JP,
                        v_align=MSO_ANCHOR.MIDDLE)

    # Vertical divider
    div_x = Inches(0.5) + col_w + Inches(0.2)
    add_rect(slide, div_x, col_hdr_y, Inches(0.02), SLIDE_H - col_hdr_y - Inches(0.7), colors["sep"])

    place_logos(slide, logos["aivalix"], logos["client"], "content")


def make_team(prs, data, colors, logos):
    """Team/company introduction slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, colors["white"])

    add_rect(slide, 0, 0, SLIDE_W, Inches(1.35), colors["black"])
    add_textbox(slide, data.get("title", "チーム紹介"),
                Inches(0.55), Inches(0.3), Inches(8), Inches(0.8),
                28, bold=True, color=colors["white"], font=FONT_JP,
                v_align=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 0, Inches(1.35), SLIDE_W, Inches(0.05), colors["accent"])

    members = data.get("members", [])
    n = len(members)
    cols = min(n, 4)
    col_w = (SLIDE_W - Inches(1.0)) / cols

    for i, member in enumerate(members):
        col = i % cols
        row = i // cols
        x = Inches(0.5) + col * col_w
        y = Inches(1.8) + row * Inches(2.5)

        # Name
        add_textbox(slide, member.get("name", ""),
                    x, y, col_w - Inches(0.3), Inches(0.5),
                    18, bold=True, color=colors["text_dark"], font=FONT_JP,
                    h_align=PP_ALIGN.CENTER, v_align=MSO_ANCHOR.MIDDLE)
        # Role
        add_textbox(slide, member.get("role", ""),
                    x, y + Inches(0.5), col_w - Inches(0.3), Inches(0.4),
                    12, color=colors["accent"], font=FONT_JP,
                    h_align=PP_ALIGN.CENTER, v_align=MSO_ANCHOR.MIDDLE)
        # Bio
        add_textbox(slide, member.get("bio", ""),
                    x, y + Inches(1.0), col_w - Inches(0.3), Inches(1.2),
                    11, color=colors["text_dark"], font=FONT_JP,
                    word_wrap=True, auto_shrink=True, min_size=9,
                    v_align=MSO_ANCHOR.TOP)

    place_logos(slide, logos["aivalix"], logos["client"], "content")


def make_back_cover(prs, data, colors, logos):
    """Thank you / closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, colors["black"])

    add_textbox(slide, data.get("title", "Thank You"),
                Inches(0.5), Inches(1.5), SLIDE_W - Inches(1.0), Inches(2.0),
                52, bold=True, color=colors["white"], font=FONT_EN,
                h_align=PP_ALIGN.CENTER, v_align=MSO_ANCHOR.MIDDLE)

    add_rect(slide, (SLIDE_W - Inches(3.0)) / 2, Inches(3.6),
             Inches(3.0), Inches(0.04), colors["accent"])

    contact = data.get("contact", "")
    add_textbox(slide, contact,
                Inches(1.0), Inches(4.0), SLIDE_W - Inches(2.0), Inches(1.5),
                14, color=colors["text_light"], font=FONT_JP,
                h_align=PP_ALIGN.CENTER, word_wrap=True,
                v_align=MSO_ANCHOR.MIDDLE)

    place_logos(slide, logos["aivalix"], logos["client"], "back_cover")


# ── Main Pipeline ───────────────────────────────────────────────────────

def build(structure_path, output_path, colors, logos, *, emit_slots=False):
    """Build presentation one slide at a time with checkpoints.

    If emit_slots=True, also writes image_slots.json next to the output
    with exact pixel dimensions for each image area.
    """
    global _IMAGE_SLOTS
    _IMAGE_SLOTS = []  # Reset for this build

    with open(structure_path, encoding="utf-8") as f:
        data = json.load(f)

    img_map = {}
    imap_path = os.path.join(os.path.dirname(structure_path), "image_map.json")
    if os.path.exists(imap_path):
        with open(imap_path, encoding="utf-8") as f:
            img_map = json.load(f)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = data.get("slides", [])
    date = data.get("date", "")
    checkpoint_dir = os.path.dirname(os.path.abspath(output_path))
    os.makedirs(checkpoint_dir, exist_ok=True)

    success_count = 0
    fail_count = 0

    for idx, sd in enumerate(slides):
        info = img_map.get(str(idx), {})
        ipath = info.get("path") if isinstance(info, dict) else info
        if ipath and not os.path.isabs(ipath):
            ipath = os.path.join(os.path.dirname(structure_path), ipath)

        stype = sd.get("type", "content")
        title_preview = sd.get("title", "")[:38]
        print(f"[{idx + 1:02}/{len(slides)}] {stype}: {title_preview}")

        try:
            if stype == "cover":
                sd["date"] = date
                make_cover(prs, sd, ipath, colors, logos, slide_idx=idx)
            elif stype == "agenda":
                make_agenda(prs, sd, colors, logos)
            elif stype == "section":
                make_section(prs, sd, ipath, colors, logos, slide_idx=idx)
            elif stype == "two_column":
                make_two_column(prs, sd, ipath, colors, logos)
            elif stype == "team":
                make_team(prs, sd, colors, logos)
            elif stype == "back_cover":
                make_back_cover(prs, sd, colors, logos)
            elif stype == "content":
                if sd.get("content_type") == "table":
                    make_content_table(prs, sd, ipath, colors, logos, slide_idx=idx)
                else:
                    make_content_bullets(prs, sd, ipath, colors, logos, slide_idx=idx)
            else:
                # Default to bullets
                make_content_bullets(prs, sd, ipath, colors, logos, slide_idx=idx)

            success_count += 1

            # Save checkpoint
            cp_path = os.path.join(checkpoint_dir, f"checkpoint_{idx + 1}.pptx")
            prs.save(cp_path)

        except Exception as e:
            fail_count += 1
            print(f"  [FAIL] {e}")
            # Add error placeholder
            try:
                err_slide = prs.slides.add_slide(prs.slide_layouts[6])
                set_bg(err_slide, colors["light_bg"])
                add_textbox(err_slide, f"[Error on slide {idx + 1}: {stype}]",
                            Inches(1), Inches(2), Inches(8), Inches(2),
                            24, bold=True, color=RGBColor(0xCC, 0x00, 0x00),
                            h_align=PP_ALIGN.CENTER)
                add_textbox(err_slide, str(e),
                            Inches(1), Inches(4), Inches(8), Inches(1),
                            12, color=colors["text_light"],
                            h_align=PP_ALIGN.CENTER)
            except Exception:
                pass

    # Save final
    prs.save(output_path)
    size_kb = os.path.getsize(output_path) // 1024
    print(f"\n[DONE] {output_path} ({size_kb} KB)")
    print(f"  OK: {success_count}, FAIL: {fail_count}")

    # Clean up checkpoints
    for idx in range(len(slides)):
        cp = os.path.join(checkpoint_dir, f"checkpoint_{idx + 1}.pptx")
        if os.path.exists(cp):
            try:
                os.remove(cp)
            except Exception:
                pass

    # Emit image slot info for two-pass workflow
    if emit_slots and _IMAGE_SLOTS:
        slots_path = os.path.join(checkpoint_dir, "image_slots.json")
        with open(slots_path, "w", encoding="utf-8") as f:
            json.dump(_IMAGE_SLOTS, f, indent=2, ensure_ascii=False)
        print(f"[SLOTS] {slots_path} — {len(_IMAGE_SLOTS)} image slots recorded")

    return output_path


if __name__ == "__main__":
    ap = argparse.ArgumentParser(description="AIVALIX Corporate Pitch PPTX Generator")
    ap.add_argument("--structure", required=True, help="Slide structure JSON")
    ap.add_argument("--output", required=True, help="Output PPTX path")
    ap.add_argument("--aivalix-logo", default="", help="AIVALIX logo path")
    ap.add_argument("--client-logo", default="", help="Client logo path")
    ap.add_argument("--client-color", default="#1A365D", help="Client accent color hex")
    ap.add_argument("--emit-image-slots", action="store_true",
                    help="Write image_slots.json with exact pixel dimensions for each image area")
    args = ap.parse_args()

    colors = dict(C)
    if args.client_color:
        colors["accent"] = hex_to_rgb(args.client_color)
        # Generate a light version of accent for use on dark backgrounds
        h = args.client_color.lstrip("#")
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        lr = min(255, r + int((255 - r) * 0.55))
        lg = min(255, g + int((255 - g) * 0.55))
        lb = min(255, b + int((255 - b) * 0.55))
        colors["accent_light"] = RGBColor(lr, lg, lb)

    logos = {
        "aivalix": args.aivalix_logo if args.aivalix_logo else None,
        "client": args.client_logo if args.client_logo else None,
    }

    build(args.structure, args.output, colors, logos,
          emit_slots=args.emit_image_slots)

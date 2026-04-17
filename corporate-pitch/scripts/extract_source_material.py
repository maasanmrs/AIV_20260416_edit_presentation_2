#!/usr/bin/env python3
"""Extract text content from PDF, DOCX, PPTX, and TXT files into structured JSON."""
import argparse, json, os, sys


def extract_pdf(path):
    """Extract text from PDF using PyPDF2."""
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        print(f"  [WARN] PyPDF2 not installed, skipping {path}")
        return None
    try:
        reader = PdfReader(path)
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                pages.append({"page": i + 1, "text": text.strip()})
        return {"filename": os.path.basename(path), "type": "pdf", "content": pages}
    except Exception as e:
        print(f"  [ERR] PDF read failed: {path} - {e}")
        return None


def extract_docx(path):
    """Extract paragraphs and tables from DOCX."""
    try:
        from docx import Document
    except ImportError:
        print(f"  [WARN] python-docx not installed, skipping {path}")
        return None
    try:
        doc = Document(path)
        paragraphs = []
        for p in doc.paragraphs:
            if p.text.strip():
                paragraphs.append({
                    "text": p.text.strip(),
                    "style": p.style.name if p.style else "Normal",
                })
        tables = []
        for t_idx, table in enumerate(doc.tables):
            rows = []
            for row in table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            if rows:
                tables.append({"table_index": t_idx, "rows": rows})
        return {
            "filename": os.path.basename(path),
            "type": "docx",
            "content": paragraphs,
            "tables": tables,
        }
    except Exception as e:
        print(f"  [ERR] DOCX read failed: {path} - {e}")
        return None


def extract_pptx(path):
    """Extract slide text and notes from PPTX."""
    try:
        from pptx import Presentation
    except ImportError:
        print(f"  [WARN] python-pptx not installed, skipping {path}")
        return None
    try:
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            slides.append({
                "slide": i + 1,
                "texts": texts,
                "notes": notes,
            })
        return {"filename": os.path.basename(path), "type": "pptx", "content": slides}
    except Exception as e:
        print(f"  [ERR] PPTX read failed: {path} - {e}")
        return None


def extract_txt(path):
    """Read plain text file."""
    try:
        with open(path, encoding="utf-8") as f:
            text = f.read().strip()
        return {"filename": os.path.basename(path), "type": "txt", "content": text}
    except Exception as e:
        print(f"  [ERR] TXT read failed: {path} - {e}")
        return None


IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".svg", ".webp", ".bmp"}


def extract_file(path):
    """Route file to the appropriate extractor."""
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return extract_pdf(path)
    elif ext in (".docx", ".doc"):
        return extract_docx(path)
    elif ext == ".pptx":
        return extract_pptx(path)
    elif ext in (".txt", ".md", ".csv"):
        return extract_txt(path)
    elif ext in IMAGE_EXTS:
        return {"filename": os.path.basename(path), "type": "image", "path": path}
    else:
        print(f"  [SKIP] Unsupported file type: {path}")
        return None


def main():
    ap = argparse.ArgumentParser(description="Extract source material from files")
    ap.add_argument("--files", nargs="+", required=True, help="Input file paths")
    ap.add_argument("--output", required=True, help="Output JSON path")
    args = ap.parse_args()

    sources = []
    images = []

    for fpath in args.files:
        fpath = fpath.strip('"').strip("'")
        if not os.path.exists(fpath):
            print(f"  [WARN] File not found: {fpath}")
            continue

        print(f"  [READ] {os.path.basename(fpath)}")
        result = extract_file(fpath)
        if result:
            if result.get("type") == "image":
                images.append(result["path"])
            else:
                sources.append(result)

    output = {"sources": sources, "extracted_images": images}
    os.makedirs(os.path.dirname(os.path.abspath(args.output)), exist_ok=True)
    with open(args.output, "w", encoding="utf-8") as f:
        json.dump(output, f, ensure_ascii=False, indent=2)
    print(f"\n[DONE] Extracted {len(sources)} source(s), {len(images)} image(s) -> {args.output}")


if __name__ == "__main__":
    main()
